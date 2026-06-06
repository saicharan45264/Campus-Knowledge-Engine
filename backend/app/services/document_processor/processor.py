"""
CurriculumLens — Document Processing Pipeline

Handles the complete document ingestion workflow:
1. Format detection (PDF, PPT, Image)
2. Text extraction
3. Page-as-image rendering (for visual embeddings)
4. Semantic text chunking
5. Embedding generation (BGE-M3 for text)
6. Vector indexing (Qdrant)
7. BM25 indexing (Elasticsearch)

This is the core data pipeline that feeds all downstream modules
(retrieval, KG construction, PYQ analysis).
"""

import io
import os
import uuid
import logging
from pathlib import Path
from typing import Optional

import fitz  # PyMuPDF
from PIL import Image

logger = logging.getLogger("curriculumlens.document_processor")


class DocumentProcessor:
    """
    Processes uploaded academic documents into indexed, retrievable chunks.

    Dual representation strategy:
    - Text path: Extract text → Chunk → Embed (BGE-M3) → Index in Qdrant
    - Visual path: Render page images → Embed (ColPali) → Index in Qdrant
    """

    def __init__(
        self,
        chunk_size: int = 512,
        chunk_overlap: int = 64,
    ):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self._text_splitter = None
        self._embedding_model = None

    @property
    def text_splitter(self):
        """Lazy-load text splitter."""
        if self._text_splitter is None:
            from langchain_text_splitters import RecursiveCharacterTextSplitter
            self._text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=self.chunk_size,
                chunk_overlap=self.chunk_overlap,
                length_function=len,
                separators=[
                    "\n\n",       # Paragraph breaks
                    "\n",         # Line breaks
                    ". ",         # Sentences
                    "; ",         # Semicolons
                    ", ",         # Commas
                    " ",          # Words
                    "",           # Characters
                ],
            )
        return self._text_splitter

    # ── PDF Processing ──

    def process_pdf(
        self,
        file_path: str,
        output_dir: str,
    ) -> dict:
        """
        Process a PDF document:
        1. Extract text per page
        2. Render each page as an image (for ColPali visual embeddings)
        3. Chunk the extracted text

        Returns: {
            "pages": [{"page_num": int, "text": str, "image_path": str}],
            "chunks": [{"text": str, "page_num": int, "chunk_index": int}],
            "page_count": int,
        }
        """
        doc = fitz.open(file_path)
        pages = []
        all_text_by_page = []

        images_dir = Path(output_dir) / "page_images"
        images_dir.mkdir(parents=True, exist_ok=True)

        for page_num in range(len(doc)):
            page = doc[page_num]

            # Extract text
            text = page.get_text("text")
            all_text_by_page.append((page_num + 1, text))

            # Render page as image (300 DPI for good quality)
            mat = fitz.Matrix(300 / 72, 300 / 72)  # 300 DPI
            pix = page.get_pixmap(matrix=mat)
            image_path = str(images_dir / f"page_{page_num + 1:04d}.png")
            pix.save(image_path)

            pages.append({
                "page_num": page_num + 1,
                "text": text.strip(),
                "image_path": image_path,
                "char_count": len(text),
            })

        doc.close()

        # Semantic chunking with page tracking
        chunks = self._chunk_pages(all_text_by_page)

        logger.info(
            "Processed PDF: %s — %d pages, %d chunks",
            file_path, len(pages), len(chunks),
        )

        return {
            "pages": pages,
            "chunks": chunks,
            "page_count": len(pages),
        }

    # ── PPT Processing ──

    def process_ppt(
        self,
        file_path: str,
        output_dir: str,
    ) -> dict:
        """
        Process a PowerPoint document:
        1. Extract text from each slide (title, body, tables, notes)
        2. Render slides as images using PyMuPDF (convert PPT → PDF first)
        3. Chunk the extracted text

        Returns same structure as process_pdf.
        """
        from pptx import Presentation

        prs = Presentation(file_path)
        pages = []
        all_text_by_page = []

        images_dir = Path(output_dir) / "slide_images"
        images_dir.mkdir(parents=True, exist_ok=True)

        for slide_num, slide in enumerate(prs.slides, 1):
            # Extract all text from slide shapes
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            texts.append(text)

                # Extract table text
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(
                            cell.text.strip() for cell in row.cells if cell.text.strip()
                        )
                        if row_text:
                            texts.append(row_text)

            # Extract slide notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    texts.append(f"[Speaker Notes] {notes}")

            slide_text = "\n".join(texts)
            all_text_by_page.append((slide_num, slide_text))

            pages.append({
                "page_num": slide_num,
                "text": slide_text,
                "image_path": None,  # PPT image rendering requires LibreOffice; skip for now
                "char_count": len(slide_text),
            })

        # Chunk
        chunks = self._chunk_pages(all_text_by_page)

        logger.info(
            "Processed PPT: %s — %d slides, %d chunks",
            file_path, len(pages), len(chunks),
        )

        return {
            "pages": pages,
            "chunks": chunks,
            "page_count": len(pages),
        }

    # ── Image Processing ──

    def process_image(
        self,
        file_path: str,
        output_dir: str,
    ) -> dict:
        """
        Process an uploaded image (handwritten notes, textbook photo, diagram).

        For images, the primary path is visual (ColPali embedding).
        Text extraction is attempted via OCR as a fallback.

        Returns same structure with single "page".
        """
        # Copy/convert image to standard format
        images_dir = Path(output_dir) / "page_images"
        images_dir.mkdir(parents=True, exist_ok=True)

        img = Image.open(file_path)
        image_path = str(images_dir / "page_0001.png")
        img.save(image_path, "PNG")

        # Text extraction would use Surya OCR here
        # For now, return with empty text (visual path is primary)
        text = ""  # TODO: Surya OCR fallback

        pages = [{
            "page_num": 1,
            "text": text,
            "image_path": image_path,
            "char_count": len(text),
        }]

        chunks = []
        if text.strip():
            chunks = self._chunk_pages([(1, text)])

        return {
            "pages": pages,
            "chunks": chunks,
            "page_count": 1,
        }

    # ── Chunking ──

    def _chunk_pages(
        self,
        pages_text: list[tuple[int, str]],
    ) -> list[dict]:
        """
        Chunk text from multiple pages while preserving page number metadata.

        Each chunk includes:
        - text: The chunk content
        - page_num: Source page number
        - chunk_index: Global chunk index within the document
        """
        chunks = []
        chunk_index = 0

        for page_num, text in pages_text:
            if not text.strip():
                continue

            page_chunks = self.text_splitter.split_text(text)

            for chunk_text in page_chunks:
                chunks.append({
                    "text": chunk_text,
                    "page_num": page_num,
                    "chunk_index": chunk_index,
                    "char_count": len(chunk_text),
                })
                chunk_index += 1

        return chunks

    # ── Dispatcher ──

    def process(
        self,
        file_path: str,
        document_type: str,
        output_dir: str,
    ) -> dict:
        """
        Route to the appropriate processor based on document type.
        """
        processors = {
            "pdf": self.process_pdf,
            "ppt": self.process_ppt,
            "image": self.process_image,
            "handwritten": self.process_image,  # Same visual path
        }

        processor = processors.get(document_type)
        if not processor:
            raise ValueError(f"Unsupported document type: {document_type}")

        return processor(file_path, output_dir)
